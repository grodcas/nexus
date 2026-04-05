#!/usr/bin/env python3
"""
parse_metadata.py — Extract metadata from OneDrive files for Nexus worktree enrichment.

Usage:
    python3 parse_metadata.py <chunk_name> <rclone_path1> [<rclone_path2> ...]

Example:
    python3 parse_metadata.py education_y1_y3 "Documentos/1º AERO" "Documentos/2º AERO" "Documentos/3º AERO"

Output: nexus/metadata/<chunk_name>.json
"""

import sys
import os
import json
import subprocess
import tempfile
import re
import time
from pathlib import Path

# Extensions to skip downloading (binary/media/archives/CAD)
SKIP_EXTENSIONS = {
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.fig', '.ico', '.svg',
    '.zip', '.rar', '.7z', '.gz', '.tar',
    '.f3d', '.f3z', '.slx', '.slxc', '.sim', '.prt', '.stl', '.dwg',
    '.bak', '.dc1', '.dwp', '.pdsprj', '.bpl', '.mat', '.dat', '.mdl',
    '.prj', '.vts', '.3ds', '.avi', '.mp4', '.wav', '.trk', '.p12',
    '.dxf', '.exe', '.lnk', '.url', '.iml', '.xml', '.json',
    '.class', '.jar', '.so', '.dll', '.dylib', '.o', '.a',
    '.woff', '.woff2', '.ttf', '.eot',
    '.pyc', '.pyo', '.egg-info',
    '.mp3', '.ogg', '.flac',
    '.tif', '.tiff', '.webp', '.heic',
    '.db', '.sqlite', '.mdb',
}

RCLONE_REMOTE = "onedrive"
DOWNLOAD_DIR = "/tmp/nexus_parse"
METADATA_DIR = os.path.expanduser("~/nexus/metadata")


def should_skip(filename):
    ext = os.path.splitext(filename)[1].lower()
    return ext in SKIP_EXTENSIONS


def build_rclone_excludes():
    """Build rclone --exclude flags for skip extensions."""
    excludes = []
    for ext in sorted(SKIP_EXTENSIONS):
        excludes.extend(["--exclude", f"*{ext}"])
    # Also skip common noise directories
    for noise in ["__pycache__", ".git", "node_modules", "site-packages", ".venv", "venv"]:
        excludes.extend(["--exclude", f"{noise}/**"])
    return excludes


def download_chunk(rclone_paths, dest_dir):
    """Download files from OneDrive via rclone, skipping binary/media."""
    os.makedirs(dest_dir, exist_ok=True)
    excludes = build_rclone_excludes()

    for rpath in rclone_paths:
        print(f"  Downloading: {RCLONE_REMOTE}:{rpath}/ ...")
        src = f"{RCLONE_REMOTE}:{rpath}"
        dst = os.path.join(dest_dir, rpath)
        os.makedirs(dst, exist_ok=True)

        cmd = ["rclone", "copy", src, dst] + excludes + [
            "--transfers", "8",
            "--checkers", "8",
            "--low-level-retries", "3",
            "--retries", "2",
            "--stats", "0",
            "-q",
        ]
        result = subprocess.run(cmd, capture_output=True, text=True)
        if result.returncode != 0:
            print(f"    Warning: rclone returned {result.returncode}: {result.stderr[:200]}")


def parse_pdf(filepath):
    """Extract metadata from PDF using PyMuPDF."""
    try:
        import fitz
        doc = fitz.open(filepath)
        meta = {
            "pages": doc.page_count,
            "title": None,
            "summary": None,
            "headings": [],
        }

        # Try metadata title
        pdf_meta = doc.metadata
        if pdf_meta and pdf_meta.get("title", "").strip():
            meta["title"] = pdf_meta["title"].strip()

        # Try TOC
        toc = doc.get_toc()
        if toc:
            meta["headings"] = [entry[1] for entry in toc[:15]]

        # Extract first page text for summary
        if doc.page_count > 0:
            text = doc[0].get_text("text").strip()
            if text:
                # Clean and truncate
                text = re.sub(r'\s+', ' ', text)
                meta["summary"] = text[:300]

        doc.close()
        return meta
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_docx(filepath):
    """Extract metadata from Word document."""
    try:
        from docx import Document
        doc = Document(filepath)

        meta = {
            "paragraphs": len(doc.paragraphs),
            "title": None,
            "summary": None,
            "headings": [],
        }

        # Core properties
        if doc.core_properties.title:
            meta["title"] = doc.core_properties.title

        # Extract headings and first text
        first_text = []
        for para in doc.paragraphs[:50]:
            if para.style and para.style.name and para.style.name.startswith('Heading'):
                meta["headings"].append(para.text.strip())
            if para.text.strip() and len(first_text) < 3:
                first_text.append(para.text.strip())

        if first_text:
            meta["summary"] = " ".join(first_text)[:300]

        return meta
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_xlsx(filepath):
    """Extract metadata from Excel spreadsheet."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, read_only=True, data_only=True)

        sheets = []
        for name in wb.sheetnames[:10]:
            ws = wb[name]
            rows = 0
            header = []
            for i, row in enumerate(ws.iter_rows(max_row=2, values_only=True)):
                if i == 0:
                    header = [str(c)[:50] for c in row if c is not None][:8]
                rows += 1
            # Count total rows (approximate for large sheets)
            try:
                rows = ws.max_row or 0
            except:
                pass
            sheets.append({"name": name, "rows": rows, "header": header})

        wb.close()
        return {"sheets": sheets}
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_xls(filepath):
    """Handle old .xls format — just report it exists."""
    return {"note": "Old .xls format — metadata extraction skipped"}


def parse_pptx(filepath):
    """Extract metadata from PowerPoint."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)

        slides = []
        for slide in prs.slides[:30]:
            title = None
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
            slides.append(title)

        return {
            "slide_count": len(prs.slides),
            "slide_titles": [s for s in slides if s],
        }
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_tex(filepath):
    """Extract metadata from LaTeX file."""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read(5000)

        meta = {"title": None, "sections": []}

        # Extract title
        m = re.search(r'\\title\{([^}]+)\}', content)
        if m:
            meta["title"] = m.group(1).strip()

        # Extract sections
        for m in re.finditer(r'\\(?:section|chapter)\{([^}]+)\}', content):
            meta["sections"].append(m.group(1).strip())

        # Extract abstract
        m = re.search(r'\\begin\{abstract\}(.*?)\\end\{abstract\}', content, re.DOTALL)
        if m:
            meta["abstract"] = re.sub(r'\s+', ' ', m.group(1).strip())[:300]

        return meta
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_code(filepath, comment_char='%'):
    """Extract first comment block and function signature from code files."""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            lines = f.readlines()[:30]

        meta = {"first_lines": [], "function_sig": None}

        # Collect leading comments
        for line in lines:
            stripped = line.strip()
            if stripped.startswith(comment_char):
                meta["first_lines"].append(stripped.lstrip(comment_char).strip())
            elif stripped.startswith('function') or stripped.startswith('def ') or stripped.startswith('class '):
                meta["function_sig"] = stripped[:100]
                break
            elif stripped == '' and not meta["first_lines"]:
                continue
            elif meta["first_lines"]:
                break

        # Build summary from comments
        if meta["first_lines"]:
            meta["summary"] = " ".join(meta["first_lines"][:5])[:300]

        return meta
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_text(filepath):
    """Read first few lines of a plain text file."""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read(1000)
        return {"summary": re.sub(r'\s+', ' ', text.strip())[:300]}
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_file(filepath):
    """Route to appropriate parser based on extension."""
    ext = os.path.splitext(filepath)[1].lower()

    if ext == '.pdf':
        return parse_pdf(filepath)
    elif ext in ('.docx', '.doc'):
        if ext == '.doc':
            return {"note": "Old .doc format — limited extraction"}
        return parse_docx(filepath)
    elif ext == '.xlsx':
        return parse_xlsx(filepath)
    elif ext == '.xls':
        return parse_xls(filepath)
    elif ext == '.pptx':
        return parse_pptx(filepath)
    elif ext == '.tex':
        return parse_tex(filepath)
    elif ext == '.m' or ext == '.asv':
        return parse_code(filepath, comment_char='%')
    elif ext == '.nb':
        return parse_text(filepath)
    elif ext in ('.py', '.mjs', '.js', '.cpp', '.cc', '.c', '.h'):
        return parse_code(filepath, comment_char='//' if ext in ('.cpp', '.cc', '.c', '.h') else '#')
    elif ext in ('.txt', '.md', '.csv', '.html', '.mhtml'):
        return parse_text(filepath)
    elif ext in ('.ipynb',):
        return parse_text(filepath)
    else:
        return None


def process_chunk(chunk_name, rclone_paths):
    """Main: download, parse, output JSON."""
    dest_dir = os.path.join(DOWNLOAD_DIR, chunk_name)

    print(f"=== Chunk: {chunk_name} ===")
    print(f"Paths: {rclone_paths}")
    print(f"Download dir: {dest_dir}")
    print()

    # Step 1: Download
    t0 = time.time()
    print("[1/3] Downloading files...")
    download_chunk(rclone_paths, dest_dir)
    t1 = time.time()
    print(f"  Download took {t1-t0:.1f}s")

    # Step 2: Parse all downloaded files
    print("[2/3] Parsing files...")
    results = []
    file_count = 0
    parse_count = 0
    error_count = 0

    for rpath in rclone_paths:
        local_root = os.path.join(dest_dir, rpath)
        if not os.path.exists(local_root):
            print(f"  Warning: {local_root} not found after download")
            continue

        for dirpath, dirnames, filenames in os.walk(local_root):
            for fname in sorted(filenames):
                file_count += 1
                fpath = os.path.join(dirpath, fname)

                # Compute OneDrive-relative path
                rel_path = os.path.relpath(fpath, dest_dir)

                if should_skip(fname):
                    continue

                meta = parse_file(fpath)
                if meta is None:
                    continue

                parse_count += 1
                if "error" in meta:
                    error_count += 1

                entry = {
                    "onedrive_path": rel_path,
                    "filename": fname,
                    "type": os.path.splitext(fname)[1].lstrip('.').lower(),
                    "size_bytes": os.path.getsize(fpath),
                }
                entry.update(meta)
                results.append(entry)

                if parse_count % 50 == 0:
                    print(f"    Parsed {parse_count} files...")

    # Step 3: Write JSON
    t2 = time.time()
    print(f"  Parsing took {t2-t1:.1f}s")
    print(f"[3/3] Writing JSON ({parse_count} entries, {error_count} errors, {file_count} total files seen)...")
    output_path = os.path.join(METADATA_DIR, f"{chunk_name}.json")
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)

    t3 = time.time()
    print(f"\nDone in {t3-t0:.1f}s! Output: {output_path}")
    print(f"  Files downloaded: {file_count}")
    print(f"  Files parsed: {parse_count}")
    print(f"  Parse errors: {error_count}")

    return results


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    chunk_name = sys.argv[1]
    rclone_paths = sys.argv[2:]

    process_chunk(chunk_name, rclone_paths)
