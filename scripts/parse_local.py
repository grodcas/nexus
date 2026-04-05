#!/usr/bin/env python3
"""
parse_local.py — Parse metadata from local OneDrive files.

Usage:
    python3 parse_local.py /tmp/nexus_all_files.txt metadata/all.json
"""

import sys
import os
import json
import re
import time

OD_PREFIX = os.path.expanduser("~/Library/CloudStorage/OneDrive-Personal/")


def parse_pdf(filepath):
    try:
        import fitz
        doc = fitz.open(filepath)
        meta = {"pages": doc.page_count, "title": None, "summary": None, "headings": []}
        pdf_meta = doc.metadata
        if pdf_meta and pdf_meta.get("title", "").strip():
            meta["title"] = pdf_meta["title"].strip()
        toc = doc.get_toc()
        if toc:
            meta["headings"] = [entry[1] for entry in toc[:15]]
        if doc.page_count > 0:
            text = doc[0].get_text("text").strip()
            if text:
                meta["summary"] = re.sub(r'\s+', ' ', text)[:300]
        doc.close()
        return meta
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_docx(filepath):
    try:
        from docx import Document
        doc = Document(filepath)
        meta = {"paragraphs": len(doc.paragraphs), "title": None, "summary": None, "headings": []}
        if doc.core_properties.title:
            meta["title"] = doc.core_properties.title
        first_text = []
        for para in doc.paragraphs[:50]:
            if para.style and para.style.name and para.style.name.startswith('Heading'):
                meta["headings"].append(para.text.strip())
            if para.text.strip() and len(first_text) < 3:
                first_text.append(para.text.strip())
        if first_text:
            meta["summary"] = " ".join(first_text)[:300]
        return meta
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_xlsx(filepath):
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, read_only=True, data_only=True)
        sheets = []
        for name in wb.sheetnames[:10]:
            ws = wb[name]
            header = []
            for i, row in enumerate(ws.iter_rows(max_row=2, values_only=True)):
                if i == 0:
                    header = [str(c)[:50] for c in row if c is not None][:8]
            try:
                rows = ws.max_row or 0
            except:
                rows = 0
            sheets.append({"name": name, "rows": rows, "header": header})
        wb.close()
        return {"sheets": sheets}
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides = []
        for slide in prs.slides[:30]:
            title = slide.shapes.title.text.strip() if slide.shapes.title else None
            slides.append(title)
        return {"slide_count": len(prs.slides), "slide_titles": [s for s in slides if s]}
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_tex(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read(5000)
        meta = {"title": None, "sections": []}
        m = re.search(r'\\title\{([^}]+)\}', content)
        if m:
            meta["title"] = m.group(1).strip()
        for m in re.finditer(r'\\(?:section|chapter)\{([^}]+)\}', content):
            meta["sections"].append(m.group(1).strip())
        m = re.search(r'\\begin\{abstract\}(.*?)\\end\{abstract\}', content, re.DOTALL)
        if m:
            meta["abstract"] = re.sub(r'\s+', ' ', m.group(1).strip())[:300]
        return meta
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_code(filepath, comment_char='%'):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            lines = f.readlines()[:30]
        meta = {"first_lines": [], "function_sig": None}
        for line in lines:
            stripped = line.strip()
            if stripped.startswith(comment_char):
                meta["first_lines"].append(stripped.lstrip(comment_char).strip())
            elif stripped.startswith('function') or stripped.startswith('def ') or stripped.startswith('class '):
                meta["function_sig"] = stripped[:100]
                break
            elif stripped == '' and not meta["first_lines"]:
                continue
            elif meta["first_lines"]:
                break
        if meta["first_lines"]:
            meta["summary"] = " ".join(meta["first_lines"][:5])[:300]
        return meta
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_text(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read(1000)
        return {"summary": re.sub(r'\s+', ' ', text.strip())[:300]}
    except Exception as e:
        return {"error": str(e)[:100]}


def parse_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == '.pdf':
        return parse_pdf(filepath)
    elif ext == '.docx':
        return parse_docx(filepath)
    elif ext == '.doc':
        return {"note": "Old .doc format"}
    elif ext == '.xlsx':
        return parse_xlsx(filepath)
    elif ext == '.xls':
        return {"note": "Old .xls format"}
    elif ext == '.pptx' or ext == '.ppt':
        if ext == '.ppt':
            return {"note": "Old .ppt format"}
        return parse_pptx(filepath)
    elif ext == '.tex':
        return parse_tex(filepath)
    elif ext in ('.m', '.asv'):
        return parse_code(filepath, '%')
    elif ext == '.nb':
        return parse_text(filepath)
    elif ext in ('.py', '.mjs', '.js', '.cpp', '.cc', '.c', '.h', '.mlx'):
        cc = '//' if ext in ('.cpp', '.cc', '.c', '.h') else '#'
        return parse_code(filepath, cc)
    elif ext in ('.txt', '.md', '.csv', '.html', '.mhtml', '.ipynb'):
        return parse_text(filepath)
    return None


def main():
    filelist = sys.argv[1]
    output = sys.argv[2]

    with open(filelist) as f:
        paths = [line.strip() for line in f if line.strip()]

    print(f"Parsing {len(paths)} files...")
    t0 = time.time()
    results = []
    errors = 0

    for i, fpath in enumerate(paths):
        if not os.path.isfile(fpath):
            continue
        meta = parse_file(fpath)
        if meta is None:
            continue
        if "error" in meta:
            errors += 1

        rel = fpath.replace(OD_PREFIX, "")
        entry = {
            "path": rel,
            "filename": os.path.basename(fpath),
            "type": os.path.splitext(fpath)[1].lstrip('.').lower(),
            "size_bytes": os.path.getsize(fpath),
        }
        entry.update(meta)
        results.append(entry)

        if (i + 1) % 2000 == 0:
            print(f"  {i+1}/{len(paths)}... ({len(results)} parsed, {errors} errors)")

    os.makedirs(os.path.dirname(output), exist_ok=True)
    # Clean surrogates from text fields
    def clean(obj):
        if isinstance(obj, str):
            return obj.encode('utf-8', errors='replace').decode('utf-8')
        if isinstance(obj, dict):
            return {k: clean(v) for k, v in obj.items()}
        if isinstance(obj, list):
            return [clean(v) for v in obj]
        return obj
    results = clean(results)
    with open(output, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=1)

    t1 = time.time()
    print(f"Done in {t1-t0:.1f}s: {len(results)} entries, {errors} errors -> {output}")


if __name__ == "__main__":
    main()
